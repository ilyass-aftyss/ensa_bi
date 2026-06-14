"""
Generate ENSA Kénitra BI — Modern PPTX Presentation
"""
import sys, os
sys.path.insert(0, "/home/runner/workspace/.pythonlibs/lib/python3.11/site-packages")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Palette ───────────────────────────────────────────────────────────────────
BLUE       = RGBColor(0x00, 0x3F, 0x8A)   # ENSA Blue
GOLD       = RGBColor(0xC8, 0x92, 0x0A)   # ENSA Gold
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK       = RGBColor(0x1A, 0x1A, 0x2E)
LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xFA)
GREEN      = RGBColor(0x1B, 0x7A, 0x4B)
RED        = RGBColor(0xC0, 0x39, 0x2B)
GRAY       = RGBColor(0x7F, 0x8C, 0x8D)
LIGHT_GRAY = RGBColor(0xF5, 0xF7, 0xFA)
ORANGE     = RGBColor(0xD3, 0x54, 0x00)
TEAL       = RGBColor(0x00, 0x7B, 0x83)

# ── Helpers ───────────────────────────────────────────────────────────────────
def W(): return Inches(13.33)
def H(): return Inches(7.5)

def solid_fill(shape, color: RGBColor):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(tf, text, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color

def add_run(para, text, size, bold=False, color=WHITE, italic=False):
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return run

def rect(slide, l, t, w, h, color, radius=0):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    solid_fill(shape, color)
    shape.line.fill.background()
    return shape

def text_box(slide, text, l, t, w, h, size=14, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = wrap
    add_text(tf, text, size, bold=bold, color=color, align=align, italic=italic)
    return txb

def bullet_box(slide, items, l, t, w, h, size=13, color=DARK, title=None, title_color=BLUE):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        add_run(p, title, size+2, bold=True, color=title_color)
    for item in items:
        p = tf.paragraphs[0] if (first and not title) else tf.add_paragraph()
        first = False
        add_run(p, f"▸  {item}", size, color=color)
    return txb

def kpi_card(slide, l, t, w, h, value, label, bg_color=BLUE, val_size=28, lbl_size=11):
    card = rect(slide, l, t, w, h, bg_color)
    # value
    txv = slide.shapes.add_textbox(Inches(l), Inches(t+0.12), Inches(w), Inches(h*0.55))
    tf = txv.text_frame
    add_text(tf, value, val_size, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # label
    txl = slide.shapes.add_textbox(Inches(l), Inches(t+0.62), Inches(w), Inches(h*0.4))
    tf2 = txl.text_frame
    add_text(tf2, label, lbl_size, color=RGBColor(0xDD,0xE8,0xFF), align=PP_ALIGN.CENTER)

def divider(slide, t, color=GOLD, thickness=0.04):
    rect(slide, 0.5, t, 12.33, thickness, color)

# ══════════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = W()
prs.slide_height = H()
blank = prs.slide_layouts[6]  # completely blank

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Cover
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)

# Full background
rect(sl, 0, 0, 13.33, 7.5, BLUE)
# Gold accent strip bottom
rect(sl, 0, 6.8, 13.33, 0.7, GOLD)
# Decorative circle top-right
circ = sl.shapes.add_shape(9, Inches(10.5), Inches(-1.2), Inches(4.5), Inches(4.5))
circ.fill.solid(); circ.fill.fore_color.rgb = RGBColor(0x00,0x55,0xAA)
circ.line.fill.background()
circ2 = sl.shapes.add_shape(9, Inches(11.2), Inches(-0.5), Inches(3.0), Inches(3.0))
circ2.fill.solid(); circ2.fill.fore_color.rgb = RGBColor(0xFF,0xFF,0xFF)
circ2.fill.fore_color._color._alpha = 0x19
circ2.line.fill.background()

# Institution
text_box(sl, "ÉCOLE NATIONALE DES SCIENCES APPLIQUÉES — KÉNITRA", 0.8, 0.5, 10, 0.5,
         size=11, color=RGBColor(0xAA,0xCC,0xFF), align=PP_ALIGN.LEFT)
text_box(sl, "Université Ibn Tofail  |  Academic Year 2024–2025", 0.8, 0.9, 10, 0.4,
         size=10, color=GRAY, align=PP_ALIGN.LEFT, italic=True)

# Gold line
rect(sl, 0.8, 1.45, 6.5, 0.05, GOLD)

# Main title
text_box(sl, "Secure Business", 0.8, 1.6, 9, 1.0, size=48, bold=True, color=WHITE)
text_box(sl, "Intelligence Platform", 0.8, 2.5, 10, 1.0, size=48, bold=True, color=WHITE)

# Subtitle
text_box(sl, "ETL Pipeline  ·  Star Schema DW  ·  JWT REST API  ·  Interactive Dashboard",
         0.8, 3.7, 11, 0.5, size=15, color=RGBColor(0xAA,0xCC,0xFF), italic=True)

# Stats row
stats = [("67,433", "Records Processed"), ("15", "API Endpoints"),
         ("6/6", "Tests Passed"), ("3", "User Roles")]
for i, (val, lbl) in enumerate(stats):
    kpi_card(sl, 0.8 + i*3.0, 4.5, 2.7, 1.5, val, lbl,
             bg_color=RGBColor(0x00,0x2A,0x5E))

# Bottom strip text
text_box(sl, "ENSA Kénitra  ·  Secure BI Platform  ·  May 2025",
         0.8, 7.0, 11, 0.4, size=10, color=BLUE, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Agenda
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, 13.33, 7.5, LIGHT_GRAY)
rect(sl, 0, 0, 13.33, 1.4, BLUE)
rect(sl, 0, 1.4, 0.08, 6.1, GOLD)

text_box(sl, "PRESENTATION AGENDA", 0.4, 0.3, 10, 0.5, size=11, color=GOLD, bold=True)
text_box(sl, "Table of Contents", 0.4, 0.7, 10, 0.6, size=26, bold=True, color=WHITE)

agenda = [
    ("01", "Project Context & Objectives",    "Why ENSA needed a BI platform"),
    ("02", "System Architecture",             "Technology stack & layers"),
    ("03", "ETL Pipeline",                    "Extract → Transform → Load"),
    ("04", "Star Schema Data Warehouse",      "Dimensional modeling & OLAP"),
    ("05", "REST API & Security",             "FastAPI + JWT + RBAC"),
    ("06", "Streamlit Dashboard",             "5 interactive analytics pages"),
    ("07", "Cloud Deployment",               "Railway + Streamlit Cloud"),
    ("08", "Results & Conclusion",            "Key metrics & future work"),
]

for i, (num, title, desc) in enumerate(agenda):
    row = i // 2
    col = i % 2
    lx = 0.5 + col * 6.3
    ty = 1.7 + row * 1.2

    bg = rect(sl, lx, ty, 5.9, 1.05, WHITE)
    rect(sl, lx, ty, 0.55, 1.05, BLUE)
    text_box(sl, num, lx+0.05, ty+0.18, 0.5, 0.6, size=18, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    text_box(sl, title, lx+0.65, ty+0.08, 5.1, 0.4, size=13, bold=True, color=BLUE)
    text_box(sl, desc,  lx+0.65, ty+0.55, 5.1, 0.4, size=10, color=GRAY, italic=True)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Context & Objectives
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, 13.33, 7.5, LIGHT_GRAY)
rect(sl, 0, 0, 13.33, 1.4, BLUE)
rect(sl, 0, 1.4, 0.08, 6.1, GOLD)

text_box(sl, "01  ·  CONTEXT & OBJECTIVES", 0.4, 0.3, 10, 0.4, size=11, color=GOLD, bold=True)
text_box(sl, "Why ENSA Kénitra Needed BI", 0.4, 0.72, 9, 0.55, size=26, bold=True, color=WHITE)

# Left — Problems
rect(sl, 0.4, 1.6, 5.9, 5.4, WHITE)
rect(sl, 0.4, 1.6, 5.9, 0.5, RED)
text_box(sl, "⚠  Challenges Before BI", 0.55, 1.65, 5.6, 0.4, size=13, bold=True, color=WHITE)

problems = [
    "No centralized visibility into academic KPIs",
    "Fragmented grade & enrollment data across files",
    "No budget execution monitoring per department",
    "Impossible to track multi-year dropout trends",
    "No unified room utilization or scheduling view",
    "Manual, error-prone Excel-based reporting",
]
for i, p in enumerate(problems):
    rect(sl, 0.55, 2.25+i*0.72, 0.28, 0.28, RED)
    text_box(sl, "✗", 0.57, 2.23+i*0.72, 0.25, 0.3, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text_box(sl, p, 0.95, 2.22+i*0.72, 5.15, 0.35, size=11, color=DARK)

# Right — Objectives
rect(sl, 6.7, 1.6, 6.2, 5.4, WHITE)
rect(sl, 6.7, 1.6, 6.2, 0.5, GREEN)
text_box(sl, "✓  Project Objectives", 6.85, 1.65, 5.8, 0.4, size=13, bold=True, color=WHITE)

objectives = [
    "Automate ETL: extract, clean & load all data",
    "Build star-schema data warehouse (OLAP-ready)",
    "Expose 15 KPI endpoints with JWT security",
    "Deploy interactive 5-page Streamlit dashboard",
    "Enforce RBAC for admin / direction / professeur",
    "Deploy on Railway + Streamlit Cloud",
]
for i, o in enumerate(objectives):
    rect(sl, 6.85, 2.25+i*0.72, 0.28, 0.28, GREEN)
    text_box(sl, "✓", 6.87, 2.23+i*0.72, 0.25, 0.3, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text_box(sl, o, 7.25, 2.22+i*0.72, 5.5, 0.35, size=11, color=DARK)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Architecture
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, 13.33, 7.5, LIGHT_GRAY)
rect(sl, 0, 0, 13.33, 1.4, BLUE)
rect(sl, 0, 1.4, 0.08, 6.1, GOLD)

text_box(sl, "02  ·  SYSTEM ARCHITECTURE", 0.4, 0.3, 10, 0.4, size=11, color=GOLD, bold=True)
text_box(sl, "5-Layer BI Architecture", 0.4, 0.72, 9, 0.55, size=26, bold=True, color=WHITE)

layers = [
    (TEAL,   "01  DATA SOURCES",   "CSV files: students, grades, finance, schedule, rooms, teachers (67,433 records)"),
    (GREEN,  "02  ETL PIPELINE",   "generate_data → extract → transform (QA) → load  |  Python + Pandas + SQLAlchemy"),
    (ORANGE, "03  DATABASE",       "PostgreSQL 14: OLTP tables (10) + Star Schema DW (7 DIM + 5 FACT) + OLAP views"),
    (RED,    "04  REST API",       "FastAPI 0.136 + uvicorn  |  JWT HS256 + RBAC  |  15 endpoints  |  Swagger UI"),
    (BLUE,   "05  DASHBOARD",      "Streamlit 1.45 + Plotly  |  5 pages  |  Hosted on Streamlit Cloud"),
]

for i, (color, title, desc) in enumerate(layers):
    ty = 1.65 + i * 1.02
    rect(sl, 0.4, ty, 12.5, 0.9, color)
    text_box(sl, title, 0.6, ty+0.05, 3.5, 0.35, size=13, bold=True, color=WHITE)
    text_box(sl, desc,  0.6, ty+0.42, 12.0, 0.35, size=10, color=RGBColor(0xDD,0xEE,0xFF))
    if i < 4:
        text_box(sl, "▼", 6.3, ty+0.92, 1.0, 0.3, size=14, color=GOLD, bold=True, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Technology Stack
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, 13.33, 7.5, LIGHT_GRAY)
rect(sl, 0, 0, 13.33, 1.4, BLUE)
rect(sl, 0, 1.4, 0.08, 6.1, GOLD)

text_box(sl, "02  ·  TECHNOLOGY STACK", 0.4, 0.3, 10, 0.4, size=11, color=GOLD, bold=True)
text_box(sl, "Complete Technology Stack", 0.4, 0.72, 9, 0.55, size=26, bold=True, color=WHITE)

categories = [
    (BLUE,  "⚙  Backend / API",
     ["Python 3.11", "FastAPI 0.136.3", "SQLAlchemy 2.0.50",
      "Pydantic 2.13.4", "python-jose (JWT)", "passlib + bcrypt 4.0.1"]),
    (GREEN, "🔄  ETL & Data",
     ["Pandas 2.2.3", "NumPy 2.2.6", "psycopg2-binary 2.9.12",
      "python-dotenv", "Custom generators", "Idempotent pipeline"]),
    (TEAL,  "📊  Dashboard",
     ["Streamlit 1.45.0", "Plotly 5.24.1", "Requests 2.34.2",
      "Multi-page app", "JWT auth client", "Cached API calls"]),
    (ORANGE,"☁  Infrastructure",
     ["Railway (API)", "Railway PostgreSQL", "Streamlit Cloud",
      "GitHub CI/CD", "Environment secrets", "HTTPS / TLS"]),
]

for i, (color, title, items) in enumerate(categories):
    col = i % 2
    row = i // 2
    lx = 0.4 + col * 6.4
    ty = 1.65 + row * 2.7

    rect(sl, lx, ty, 6.0, 2.5, WHITE)
    rect(sl, lx, ty, 6.0, 0.45, color)
    text_box(sl, title, lx+0.15, ty+0.06, 5.7, 0.35, size=13, bold=True, color=WHITE)

    for j, item in enumerate(items):
        c = j % 2; r2 = j // 2
        ix = lx + 0.2 + c * 2.9
        iy = ty + 0.58 + r2 * 0.6
        rect(sl, ix, iy, 2.6, 0.45, RGBColor(0xF0,0xF4,0xF8))
        text_box(sl, item, ix+0.1, iy+0.07, 2.4, 0.3, size=10, color=DARK, bold=(j==0))

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — ETL Pipeline
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, 13.33, 7.5, LIGHT_GRAY)
rect(sl, 0, 0, 13.33, 1.4, BLUE)
rect(sl, 0, 1.4, 0.08, 6.1, GOLD)

text_box(sl, "03  ·  ETL PIPELINE", 0.4, 0.3, 10, 0.4, size=11, color=GOLD, bold=True)
text_box(sl, "Extract → Transform → Load", 0.4, 0.72, 9, 0.55, size=26, bold=True, color=WHITE)

phases = [
    (TEAL,   "GENERATE",  "generate_data.py",
     ["1,200 students", "59,540 grades", "54 teachers", "3,000 schedules", "67,433 total records"]),
    (GREEN,  "EXTRACT",   "extract.py",
     ["Reads 10 CSV files", "Type casting", "Date parsing", "DataFrames created", "Row count logged"]),
    (ORANGE, "TRANSFORM", "transform.py",
     ["Null imputation", "Anomaly detection", "Email repair", "Grade validation", "Finance QA flags"]),
    (BLUE,   "LOAD",      "load.py",
     ["Auto-creates tables", "TRUNCATE + INSERT", "Idempotent runs", "500-row chunks", "~10s total"]),
]

for i, (color, title, module, items) in enumerate(phases):
    lx = 0.3 + i * 3.2
    rect(sl, lx, 1.7, 3.0, 5.4, WHITE)
    rect(sl, lx, 1.7, 3.0, 0.8, color)
    text_box(sl, title, lx+0.1, 1.76, 2.8, 0.35, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text_box(sl, module, lx+0.1, 2.12, 2.8, 0.28, size=9, italic=True, color=RGBColor(0xDD,0xEE,0xFF), align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        text_box(sl, f"▸  {item}", lx+0.15, 2.65+j*0.6, 2.7, 0.5, size=11, color=DARK)

    if i < 3:
        text_box(sl, "→", 3.15+i*3.2, 4.25, 0.5, 0.5, size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# Quality metrics banner
rect(sl, 0.3, 7.05, 12.7, 0.35, RGBColor(0xE8,0xF8,0xEE))
text_box(sl, "✓  0 invalid grades  |  2 salaries imputed  |  2 finance anomalies flagged  |  1,200 emails validated",
         0.5, 7.07, 12.3, 0.28, size=10, color=GREEN, align=PP_ALIGN.CENTER, bold=True)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Data Warehouse
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, 13.33, 7.5, LIGHT_GRAY)
rect(sl, 0, 0, 13.33, 1.4, BLUE)
rect(sl, 0, 1.4, 0.08, 6.1, GOLD)

text_box(sl, "04  ·  STAR SCHEMA DATA WAREHOUSE", 0.4, 0.3, 10, 0.4, size=11, color=GOLD, bold=True)
text_box(sl, "Dimensional Modeling — Kimball Methodology", 0.4, 0.72, 9, 0.55, size=24, bold=True, color=WHITE)

# Central fact box
rect(sl, 4.9, 3.0, 3.5, 1.35, BLUE)
text_box(sl, "FACT TABLES", 4.95, 3.08, 3.4, 0.35, size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
for i, f in enumerate(["fact_grades", "fact_enrollments", "fact_finance"]):
    text_box(sl, f"▸  {f}", 5.05, 3.48+i*0.27, 3.2, 0.25, size=9.5, color=WHITE)

# Dimension cards
dims = [
    (1.0, 1.7, "dim_student",    "1,200 students\nFilière, year, origin"),
    (4.9, 1.7, "dim_time",       "10 periods\n5 years × 2 semesters"),
    (8.8, 1.7, "dim_course",     "87 courses\nCredits, filière, semester"),
    (1.0, 5.0, "dim_teacher",    "54 teachers\nGrade, dept, salary"),
    (4.9, 5.0, "dim_room",       "20 rooms\nAmphis, TP, Labs"),
    (8.8, 5.0, "dim_filiere",    "5 tracks\nGI-BD, GI-GL, GSTR…"),
]

for lx, ty, name, desc in dims:
    rect(sl, lx, ty, 3.0, 1.2, WHITE)
    rect(sl, lx, ty, 3.0, 0.38, GOLD)
    text_box(sl, name, lx+0.1, ty+0.04, 2.8, 0.3, size=11, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    text_box(sl, desc, lx+0.1, ty+0.45, 2.8, 0.65, size=9.5, color=DARK, align=PP_ALIGN.CENTER)

# OLAP views note
rect(sl, 0.4, 6.55, 12.5, 0.75, RGBColor(0xE8,0xF0,0xFA))
text_box(sl, "5 Pre-built OLAP Views:", 0.6, 6.6, 2.5, 0.3, size=10, bold=True, color=BLUE)
text_box(sl,
    "olap_performance_par_filiere  ·  olap_inscriptions  ·  olap_budget_execution  ·  olap_occupation_salles  ·  olap_charge_enseignants",
    3.0, 6.6, 10.0, 0.3, size=9.5, color=DARK)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — API & Security
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, 13.33, 7.5, LIGHT_GRAY)
rect(sl, 0, 0, 13.33, 1.4, BLUE)
rect(sl, 0, 1.4, 0.08, 6.1, GOLD)

text_box(sl, "05  ·  REST API & SECURITY", 0.4, 0.3, 10, 0.4, size=11, color=GOLD, bold=True)
text_box(sl, "FastAPI + JWT HS256 + RBAC", 0.4, 0.72, 9, 0.55, size=26, bold=True, color=WHITE)

# JWT Flow
rect(sl, 0.4, 1.6, 7.8, 2.2, WHITE)
rect(sl, 0.4, 1.6, 7.8, 0.42, BLUE)
text_box(sl, "🔐  JWT Authentication Flow", 0.6, 1.66, 7.5, 0.3, size=12, bold=True, color=WHITE)

steps = ["POST /auth/token\n(username+password)", "Verify bcrypt\nhash", "Generate JWT\n(HS256, 60min)", "Return\naccess_token", "Bearer token\non all requests"]
colors_jwt = [TEAL, ORANGE, GREEN, BLUE, RED]
for i, (step, c) in enumerate(zip(steps, colors_jwt)):
    rect(sl, 0.55+i*1.48, 2.15, 1.3, 1.25, c)
    text_box(sl, step, 0.57+i*1.48, 2.22, 1.26, 1.1, size=8.5, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        text_box(sl, "→", 1.75+i*1.48, 2.6, 0.3, 0.35, size=14, color=GOLD, bold=True, align=PP_ALIGN.CENTER)

# RBAC table
rect(sl, 8.6, 1.6, 4.33, 2.2, WHITE)
rect(sl, 8.6, 1.6, 4.33, 0.42, BLUE)
text_box(sl, "👥  RBAC Permission Matrix", 8.75, 1.66, 4.0, 0.3, size=12, bold=True, color=WHITE)

headers = ["Permission", "Admin", "Dir.", "Prof."]
widths  = [2.0, 0.65, 0.65, 0.65]
lxs = [8.65, 10.65, 11.3, 11.95]
for j, (h, lx) in enumerate(zip(headers, lxs)):
    text_box(sl, h, lx, 2.1, widths[j], 0.28, size=9.5, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

perms = [
    ("Academic KPIs",    "✓", "✓", "✓"),
    ("Financial KPIs",   "✓", "✓", "✗"),
    ("HR Data",          "✓", "✓", "✗"),
    ("Write Access",     "✓", "✗", "✗"),
]
for i, (perm, a, d, p) in enumerate(perms):
    bg = RGBColor(0xF5,0xF7,0xFA) if i % 2 == 0 else WHITE
    rect(sl, 8.65, 2.45+i*0.33, 4.25, 0.3, bg)
    text_box(sl, perm, 8.68, 2.47+i*0.33, 2.0, 0.27, size=9, color=DARK)
    for j, (val, lx) in enumerate(zip([a, d, p], [10.65, 11.3, 11.95])):
        c = GREEN if val == "✓" else RED
        text_box(sl, val, lx, 2.47+i*0.33, 0.6, 0.27, size=11, color=c, align=PP_ALIGN.CENTER, bold=True)

# Endpoints grid
rect(sl, 0.4, 3.95, 12.5, 0.38, BLUE)
text_box(sl, "API Endpoint Catalog — 15 Endpoints", 0.6, 4.0, 12.0, 0.28, size=12, bold=True, color=WHITE)

endpoints = [
    ("POST", "/auth/token",           "Authenticate → JWT"),
    ("GET",  "/students/",            "Paginated student list"),
    ("GET",  "/students/by-filiere",  "Count per track"),
    ("GET",  "/students/evolution",   "Enrollment trends"),
    ("GET",  "/kpi/overview",         "Dashboard aggregates"),
    ("GET",  "/kpi/success-rate",     "Pass rate by filière"),
    ("GET",  "/kpi/abandon-rate",     "Dropout analysis"),
    ("GET",  "/kpi/avg-grade",        "Average grades"),
    ("GET",  "/kpi/top-filieres",     "Filière ranking"),
    ("GET",  "/kpi/teacher-workload", "Teaching hours"),
    ("GET",  "/kpi/room-occupancy",   "Room utilization"),
    ("GET",  "/kpi/enrollments-summary", "Enrollment stats"),
]
for i, (method, path, desc) in enumerate(endpoints):
    col = i % 3
    row = i // 3
    lx = 0.4 + col * 4.3
    ty = 4.42 + row * 0.7
    mc = GREEN if method == "GET" else ORANGE
    rect(sl, lx, ty, 0.55, 0.52, mc)
    text_box(sl, method, lx+0.02, ty+0.1, 0.5, 0.3, size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(sl, lx+0.57, ty, 3.65, 0.52, WHITE)
    text_box(sl, path, lx+0.65, ty+0.04, 2.2, 0.22, size=8.5, bold=True, color=BLUE)
    text_box(sl, desc, lx+0.65, ty+0.27, 2.9, 0.22, size=8, color=GRAY, italic=True)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Dashboard
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, 13.33, 7.5, LIGHT_GRAY)
rect(sl, 0, 0, 13.33, 1.4, BLUE)
rect(sl, 0, 1.4, 0.08, 6.1, GOLD)

text_box(sl, "06  ·  STREAMLIT DASHBOARD", 0.4, 0.3, 10, 0.4, size=11, color=GOLD, bold=True)
text_box(sl, "5-Page Interactive Analytics", 0.4, 0.72, 9, 0.55, size=26, bold=True, color=WHITE)

pages_info = [
    (BLUE,   "🏠",  "Overview",          "4 KPI cards: enrollment, success rate,\ndropout rate, budget execution"),
    (GREEN,  "📈",  "Performance",       "Grade distributions, pass/fail trends,\nfilière rankings, year comparisons"),
    (TEAL,   "🏫",  "Resources",         "Room occupancy heatmap, teaching loads,\ncredit distribution per track"),
    (ORANGE, "💰",  "Finance",           "Budget vs. expenditure charts,\nanomaly alerts, execution trends"),
    (RED,    "👥",  "Human Resources",   "Enrollment evolution, faculty workload,\ndepartment distributions"),
]

for i, (color, icon, title, desc) in enumerate(pages_info):
    col = i % 3 if i < 3 else (i - 3)
    row = 0 if i < 3 else 1
    lx = 0.4 + (col if i < 3 else col + 1.65) * 4.0
    ty = 1.7 + row * 2.8

    if i == 3:
        lx = 0.4 + 0.0 * 4.0 + 2.0
    elif i == 4:
        lx = 0.4 + 1.0 * 4.0 + 2.0

    rect(sl, lx, ty, 3.7, 2.45, WHITE)
    rect(sl, lx, ty, 3.7, 0.55, color)
    text_box(sl, f"{icon}  {title}", lx+0.15, ty+0.1, 3.4, 0.35, size=14, bold=True, color=WHITE)
    text_box(sl, desc, lx+0.15, ty+0.72, 3.4, 1.5, size=10.5, color=DARK)

# Simple layout for 5 cards: 3 top + 2 bottom centered
# Redo cleaner
for shape in list(sl.shapes)[7:]:
    sp = shape._element
    sp.getparent().remove(sp)

cols3 = [0.4, 4.5, 8.6]
cols2 = [2.45, 6.55]

for i, (color, icon, title, desc) in enumerate(pages_info[:3]):
    lx = cols3[i]
    ty = 1.7
    rect(sl, lx, ty, 3.8, 2.5, WHITE)
    rect(sl, lx, ty, 3.8, 0.55, color)
    text_box(sl, f"{icon}  {title}", lx+0.15, ty+0.1, 3.5, 0.35, size=14, bold=True, color=WHITE)
    text_box(sl, desc, lx+0.15, ty+0.72, 3.5, 1.55, size=10.5, color=DARK)

for i, (color, icon, title, desc) in enumerate(pages_info[3:]):
    lx = cols2[i]
    ty = 4.5
    rect(sl, lx, ty, 3.8, 2.5, WHITE)
    rect(sl, lx, ty, 3.8, 0.55, color)
    text_box(sl, f"{icon}  {title}", lx+0.15, ty+0.1, 3.5, 0.35, size=14, bold=True, color=WHITE)
    text_box(sl, desc, lx+0.15, ty+0.72, 3.5, 1.55, size=10.5, color=DARK)

text_box(sl, "All pages authenticate via JWT · Streamlit @st.cache_data (300s TTL) for API response caching",
         0.5, 7.12, 12.3, 0.3, size=9.5, color=GRAY, italic=True, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Deployment
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, 13.33, 7.5, LIGHT_GRAY)
rect(sl, 0, 0, 13.33, 1.4, BLUE)
rect(sl, 0, 1.4, 0.08, 6.1, GOLD)

text_box(sl, "07  ·  CLOUD DEPLOYMENT", 0.4, 0.3, 10, 0.4, size=11, color=GOLD, bold=True)
text_box(sl, "Railway + Streamlit Cloud", 0.4, 0.72, 9, 0.55, size=26, bold=True, color=WHITE)

# Railway block
rect(sl, 0.4, 1.65, 6.0, 5.3, WHITE)
rect(sl, 0.4, 1.65, 6.0, 0.5, RGBColor(0x7B, 0x00, 0xD4))
text_box(sl, "🚀  Railway (Backend)", 0.6, 1.72, 5.6, 0.35, size=14, bold=True, color=WHITE)

railway_items = [
    ("Service", "FastAPI + uvicorn"),
    ("DB",      "PostgreSQL 14+ (managed)"),
    ("Start",   "uvicorn app.main:app --host 0.0.0.0 --port $PORT"),
    ("Root Dir","bi-project/"),
    ("Auto-var","DATABASE_URL (injected)"),
    ("Manual",  "SECRET_KEY, ENVIRONMENT"),
]
for i, (k, v) in enumerate(railway_items):
    rect(sl, 0.55, 2.28+i*0.72, 1.1, 0.55, RGBColor(0x7B,0x00,0xD4))
    text_box(sl, k, 0.57, 2.34+i*0.72, 1.05, 0.35, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text_box(sl, v, 1.72, 2.34+i*0.72, 4.5, 0.35, size=10, color=DARK)

# Streamlit Cloud block
rect(sl, 6.8, 1.65, 6.1, 5.3, WHITE)
rect(sl, 6.8, 1.65, 6.1, 0.5, RGBColor(0xFF, 0x4B, 0x4B))
text_box(sl, "🎈  Streamlit Cloud (Frontend)", 7.0, 1.72, 5.7, 0.35, size=14, bold=True, color=WHITE)

sc_items = [
    ("Repo",    "github.com/<user>/<repo>"),
    ("Branch",  "main"),
    ("Main file","bi-project/dashboard/streamlit_app.py"),
    ("Reqs",    "bi-project/requirements-dashboard.txt"),
    ("Secret",  "API_URL = https://api.up.railway.app"),
    ("Python",  "3.11 (auto-detected)"),
]
for i, (k, v) in enumerate(sc_items):
    rect(sl, 6.95, 2.28+i*0.72, 1.1, 0.55, RGBColor(0xFF,0x4B,0x4B))
    text_box(sl, k, 6.97, 2.34+i*0.72, 1.05, 0.35, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text_box(sl, v, 8.12, 2.34+i*0.72, 4.6, 0.35, size=10, color=DARK)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — Results & Validation
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, 13.33, 7.5, LIGHT_GRAY)
rect(sl, 0, 0, 13.33, 1.4, BLUE)
rect(sl, 0, 1.4, 0.08, 6.1, GOLD)

text_box(sl, "08  ·  RESULTS & VALIDATION", 0.4, 0.3, 10, 0.4, size=11, color=GOLD, bold=True)
text_box(sl, "Testing, Metrics & Key Results", 0.4, 0.72, 9, 0.55, size=26, bold=True, color=WHITE)

# KPI row
kpis = [
    (BLUE,  "6/6",     "Tests Passed\n100% Coverage"),
    (GREEN, "67,433",  "Records\nProcessed"),
    (TEAL,  "<13s",    "ETL Pipeline\nDuration"),
    (GOLD,  "12",      "DW Tables\n(7 DIM + 5 FACT)"),
    (RED,   "15",      "Secured API\nEndpoints"),
]
for i, (color, val, lbl) in enumerate(kpis):
    kpi_card(sl, 0.4+i*2.56, 1.65, 2.35, 1.5, val, lbl, bg_color=color)

# Test suite table
rect(sl, 0.4, 3.3, 7.5, 0.4, BLUE)
text_box(sl, "Validation Test Suite", 0.6, 3.35, 7.0, 0.28, size=12, bold=True, color=WHITE)

tests = [
    ("FastAPI & Security Imports", "All modules importable, no missing deps"),
    ("JWT + RBAC Engine",          "Token generation, role enforcement verified"),
    ("ETL Extract Phase",          "10 CSV files read, row counts validated"),
    ("ETL Transform Phase",        "Anomalies detected, nulls imputed"),
    ("Star Schema SQL",            "DDL valid, all 12 tables + 5 views"),
    ("Dashboard Utils",            "Constants, colors, API client structure"),
]
for i, (test, detail) in enumerate(tests):
    bg = RGBColor(0xF5,0xF7,0xFA) if i % 2 == 0 else WHITE
    rect(sl, 0.4, 3.75+i*0.5, 7.5, 0.47, bg)
    rect(sl, 0.4, 3.75+i*0.5, 0.5, 0.47, GREEN)
    text_box(sl, "✓", 0.41, 3.82+i*0.5, 0.45, 0.3, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text_box(sl, test, 1.0, 3.82+i*0.5, 2.8, 0.3, size=10, bold=True, color=DARK)
    text_box(sl, detail, 3.9, 3.82+i*0.5, 3.9, 0.3, size=9.5, color=GRAY, italic=True)

# Right: Security checklist
rect(sl, 8.2, 3.3, 4.73, 3.7, WHITE)
rect(sl, 8.2, 3.3, 4.73, 0.4, BLUE)
text_box(sl, "Security Controls", 8.4, 3.35, 4.3, 0.28, size=12, bold=True, color=WHITE)

security = [
    "bcrypt password hashing",
    "JWT HS256 token signing",
    "Token expiry (60 min)",
    "Role-based route guards",
    "CORS headers configured",
    "Env-var secret isolation",
    "HTTPS via Railway proxy",
    "Law 09-08 compliance",
]
for i, s in enumerate(security):
    rect(sl, 8.3, 3.78+i*0.38, 0.32, 0.3, GREEN)
    text_box(sl, "✓", 8.31, 3.79+i*0.38, 0.3, 0.25, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text_box(sl, s, 8.7, 3.79+i*0.38, 4.1, 0.28, size=10, color=DARK)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — Conclusion
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, 13.33, 7.5, BLUE)
rect(sl, 0, 6.5, 13.33, 1.0, GOLD)

# Decorative
circ = sl.shapes.add_shape(9, Inches(-1.5), Inches(4.5), Inches(5.0), Inches(5.0))
circ.fill.solid(); circ.fill.fore_color.rgb = RGBColor(0x00,0x55,0xAA)
circ.line.fill.background()

text_box(sl, "08  ·  CONCLUSION", 0.6, 0.5, 10, 0.4, size=11, color=GOLD, bold=True)
text_box(sl, "Key Achievements", 0.6, 0.95, 9, 0.6, size=32, bold=True, color=WHITE)
rect(sl, 0.6, 1.62, 8.0, 0.06, GOLD)

achievements = [
    "✔  Complete ETL pipeline — 67,433 records in <13 seconds",
    "✔  Star-schema DW — 12 tables, 5 OLAP views, sub-second queries",
    "✔  15 secured API endpoints — JWT HS256 + 3-role RBAC",
    "✔  5-page interactive dashboard — 4 BI decision domains",
    "✔  Cloud-ready — Railway (API) + Streamlit Cloud (dashboard)",
    "✔  6/6 tests passed — 100% validation coverage",
]
for i, a in enumerate(achievements):
    text_box(sl, a, 0.7, 1.82+i*0.7, 8.5, 0.55, size=13, color=WHITE)

# Future work
rect(sl, 9.2, 1.55, 3.7, 4.6, RGBColor(0x00,0x2A,0x5E))
text_box(sl, "Future Work", 9.35, 1.65, 3.4, 0.35, size=13, bold=True, color=GOLD)
future = [
    "📡  Real data connectors\n    (ERP / SIS integration)",
    "⏰  Scheduled ETL\n    (Airflow cron jobs)",
    "🤖  Predictive ML\n    (dropout risk model)",
    "📱  Mobile dashboard\n    (React Native app)",
    "🔔  Alerting system\n    (budget anomaly SMS)",
]
for i, f in enumerate(future):
    text_box(sl, f, 9.35, 2.1+i*0.77, 3.4, 0.65, size=9.5, color=WHITE)

# Bottom bar
text_box(sl, "Thank you  ·  Questions?  ·  ENSA Kénitra BI Platform  ·  May 2025",
         0.5, 6.6, 12.3, 0.35, size=13, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(__file__), "ENSA_BI_Presentation.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
